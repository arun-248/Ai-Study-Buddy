import os
import json
import requests
import numpy as np
import re
from dotenv import load_dotenv
from datetime import datetime
from typing import List, Optional
from fastapi import FastAPI, File, UploadFile, HTTPException
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel
from PyPDF2 import PdfReader
from pptx import Presentation
from sentence_transformers import SentenceTransformer

# -----------------------
# Configuration & Setup
# -----------------------
app = FastAPI(title="AI Study Buddy Backend (Final Stable Build)")

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],  # change this in production
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)
load_dotenv()
GEMINI_API_KEY = os.getenv("GEMINI_API_KEY")
GEMINI_TEXT_URL = (
    f"https://generativelanguage.googleapis.com/v1/models/gemini-2.5-flash:generateContent?key={GEMINI_API_KEY}"
    if GEMINI_API_KEY
    else None
)

# Store documents and feedbacks in memory
DOCUMENTS = {}
FEEDBACKS = []

# Local sentence transformer for fallback
EMBED_MODEL_NAME = "all-MiniLM-L6-v2"
embedder = SentenceTransformer(EMBED_MODEL_NAME)

# -----------------------
# Helper Functions
# -----------------------
def extract_text(file: UploadFile) -> str:
    """Extract text from PDF, PPTX, or TXT"""
    name = file.filename.lower()
    file.file.seek(0)
    if name.endswith(".pdf"):
        try:
            reader = PdfReader(file.file)
            return "\n".join([page.extract_text() or "" for page in reader.pages]).strip()
        except Exception as e:
            raise HTTPException(status_code=400, detail=f"PDF parsing error: {e}")
    elif name.endswith(".pptx") or name.endswith(".ppt"):
        try:
            prs = Presentation(file.file)
            return "\n".join(
                [shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")]
            ).strip()
        except Exception as e:
            raise HTTPException(status_code=400, detail=f"PPTX parsing error: {e}")
    else:
        try:
            content = file.file.read()
            return content.decode("utf-8", errors="ignore").strip()
        except Exception as e:
            raise HTTPException(status_code=400, detail=f"Text read error: {e}")


def split_sentences(text: str) -> List[str]:
    """Split text into sentences for local summarization"""
    sents, buf = [], ""
    for ch in text:
        buf += ch
        if ch in ".!?":
            if len(buf.strip()) > 20:
                sents.append(buf.strip())
            buf = ""
    if buf.strip():
        sents.append(buf.strip())
    return sents[:1000]


def extractive_summary(text: str, n: int = 5) -> str:
    """Local extractive summarizer using embeddings"""
    sents = split_sentences(text)
    if len(sents) <= n:
        return "\n".join(sents)
    emb = embedder.encode(sents, convert_to_numpy=True)
    centroid = np.mean(emb, axis=0)
    sims = np.dot(emb, centroid) / (np.linalg.norm(emb, axis=1) * np.linalg.norm(centroid) + 1e-9)
    top = np.argsort(-sims)[:n]
    return "\n".join([sents[i] for i in sorted(top)])


def call_gemini(prompt: str) -> str:
    """Call Gemini API for text generation"""
    if not GEMINI_TEXT_URL:
        raise RuntimeError("Gemini not configured")
    headers = {"Content-Type": "application/json"}
    body = {"contents": [{"parts": [{"text": prompt}]}]}
    r = requests.post(GEMINI_TEXT_URL, headers=headers, data=json.dumps(body), timeout=40)
    r.raise_for_status()
    data = r.json()
    try:
        return data["candidates"][0]["content"]["parts"][0]["text"]
    except Exception:
        return json.dumps(data)


def get_style_specific_prompt(style_key: str, context: str) -> str:
    """Generate highly specific prompts for each summary style"""
    
    if style_key == "simple":
        return f"""You are a study assistant creating a SIMPLE SUMMARY.

STRICT REQUIREMENTS:
- Create EXACTLY 5 short bullet points
- Each bullet point should be 1-2 sentences maximum
- Focus on the most important concepts only
- Use simple, clear language
- Start each bullet with a relevant emoji

Context to summarize:
{context}

Format your response as:
• [Emoji] [Short key point 1]
• [Emoji] [Short key point 2]
• [Emoji] [Short key point 3]
• [Emoji] [Short key point 4]
• [Emoji] [Short key point 5]"""

    elif style_key == "detailed":
        return f"""You are a study assistant creating a DETAILED COMPREHENSIVE SUMMARY.

STRICT REQUIREMENTS:
- Provide in-depth explanation with 3-4 main sections
- Each section should have:
  * A clear heading
  * 2-3 paragraphs of detailed explanation
  * Examples or definitions where relevant
- Total length should be 300-400 words
- Use markdown headers (##) for sections
- Include technical terms with explanations

Context to summarize:
{context}

Format your response with clear sections and detailed explanations."""

    elif style_key == "concept":
        return f"""You are a study assistant creating a HIERARCHICAL CONCEPT MAP.

STRICT REQUIREMENTS:
- Create a multi-level bullet hierarchy showing topic relationships
- Use proper indentation to show parent-child relationships
- Include 1 main concept at the top
- 3-4 major subconcepts (indented once)
- 2-3 supporting details under each subconcept (indented twice)
- Use arrows (→) or dashes (-) to show connections
- Keep it visual and structured

Context to summarize:
{context}

Format as a hierarchical structure:
# Main Concept
  - Major Subconcept 1
    → Supporting detail
    → Supporting detail
  - Major Subconcept 2
    → Supporting detail
    → Supporting detail"""

    elif style_key == "qa":
        return f"""You are a study assistant creating a Q&A SUMMARY.

STRICT REQUIREMENTS:
- Generate EXACTLY 5 important questions
- Each question must:
  * Be specific and focused
  * Start with What/Why/How/Explain
  * Have a clear, comprehensive answer (2-4 sentences)
- Format each Q&A clearly with numbering
- Cover different aspects of the material

Context to summarize:
{context}

Format your response as:
**Q1: [Question]**
A: [Detailed answer]

**Q2: [Question]**
A: [Detailed answer]

(Continue for all 5 questions)"""

    elif style_key == "takeaways":
        return f"""You are a study assistant creating KEY TAKEAWAYS.

STRICT REQUIREMENTS:
- List EXACTLY 10 essential learning points
- Each takeaway should be:
  * Actionable or memorable
  * 1-2 sentences long
  * Numbered clearly
  * Focused on what students should remember
- Mix concepts, facts, and practical insights
- Use bold text for emphasis on key terms

Context to summarize:
{context}

Format your response as:
**Key Takeaways:**

1. [Important takeaway with **key term** bolded]
2. [Important takeaway with **key term** bolded]
...
10. [Important takeaway with **key term** bolded]"""

    else:
        # Default fallback
        return f"""Summarize the following content clearly and concisely:

{context}"""


def generate_local_fallback(style_key: str, text: str) -> str:
    """Generate style-specific local fallback summaries"""
    sents = split_sentences(text)
    
    if style_key == "simple":
        summary_sents = sents[:5] if len(sents) >= 5 else sents
        return "**Simple Summary:**\n\n" + "\n".join([f"• {s}" for s in summary_sents])
    
    elif style_key == "detailed":
        summary = extractive_summary(text, 10)
        return f"**Detailed Summary:**\n\n{summary}\n\n*Note: This is a local fallback. For better results, configure Gemini API.*"
    
    elif style_key == "concept":
        key_sents = sents[:8] if len(sents) >= 8 else sents
        concept_map = "**Concept Map:**\n\n# Main Topic\n"
        for i, s in enumerate(key_sents):
            if i % 2 == 0:
                concept_map += f"  - {s[:80]}...\n"
            else:
                concept_map += f"    → {s[:60]}...\n"
        return concept_map
    
    elif style_key == "qa":
        qa_list = "**Q&A Summary:**\n\n"
        for i, s in enumerate(sents[:5]):
            qa_list += f"**Q{i+1}: What about {s[:40]}...?**\n"
            qa_list += f"A: {s}\n\n"
        return qa_list
    
    elif style_key == "takeaways":
        key_sents = sents[:10] if len(sents) >= 10 else sents
        return "**Key Takeaways:**\n\n" + "\n".join([f"{i+1}. {s}" for i, s in enumerate(key_sents)])
    
    else:
        return extractive_summary(text, 6)


def generate_concise_study_plan(topic: str, days: int) -> str:
    """Generate a clean, concise study plan"""
    
    plan = f"📚 {days}-Day Study Plan: {topic}\n"
    plan += f"{'─' * 50}\n\n"
    
    # Divide days into phases
    fundamentals = max(1, int(days * 0.3))
    intermediate = max(1, int(days * 0.4))
    advanced = days - fundamentals - intermediate
    
    current_day = 1
    
    # Phase 1: Fundamentals
    plan += "PHASE 1: FUNDAMENTALS\n\n"
    for i in range(fundamentals):
        plan += f"📅 Day {current_day}: Core Concepts\n"
        plan += f"• Learn basic terminology and definitions\n"
        plan += f"• Understand foundational principles\n"
        plan += f"• Study real-world applications\n"
        plan += f"• Practice: Complete 3-5 beginner exercises\n"
        plan += f"⏱ Time: 2-3 hours\n\n"
        current_day += 1
    
    # Phase 2: Intermediate
    plan += "PHASE 2: BUILDING SKILLS\n\n"
    for i in range(intermediate):
        plan += f"📅 Day {current_day}: Intermediate Topics\n"
        plan += f"• Apply concepts to practical problems\n"
        plan += f"• Work through guided examples\n"
        plan += f"• Understand common patterns\n"
        plan += f"• Practice: Build a small project\n"
        plan += f"⏱ Time: 2-3 hours\n\n"
        current_day += 1
        
        # Add quiz checkpoint every 3 days
        if current_day % 3 == 0 and current_day <= days:
            plan += f"🧪 Checkpoint Quiz: Test Days {current_day-2}-{current_day-1}\n\n"
    
    # Phase 3: Advanced
    if advanced > 0:
        plan += "PHASE 3: ADVANCED PRACTICE\n\n"
        for i in range(advanced):
            plan += f"📅 Day {current_day}: Advanced Techniques\n"
            plan += f"• Master complex concepts\n"
            plan += f"• Integrate multiple topics\n"
            plan += f"• Work on real-world scenarios\n"
            plan += f"• Practice: Complete a comprehensive project\n"
            plan += f"⏱ Time: 3-4 hours\n\n"
            current_day += 1
    
    # Study tips
    plan += f"{'─' * 50}\n"
    plan += "💡 QUICK STUDY TIPS\n\n"
    plan += "1. 🍅 Use Pomodoro: 25 min focus, 5 min break\n"
    plan += "2. 📝 Take notes in your own words\n"
    plan += "3. 🔄 Review previous day before starting\n"
    plan += "4. 💻 Practice > Theory - build projects\n"
    plan += "5. 👥 Join study groups or forums\n"
    plan += "6. 😴 Get good sleep for retention\n"
    plan += "7. 🎯 Set daily goals and track progress\n\n"
    
    plan += f"{'─' * 50}\n"
    plan += f"🎯 Goal: Build strong foundation in {topic}\n"
    plan += f"📊 Track your progress daily!\n"
    plan += f"🚀 You've got this!\n"
    
    return plan


# -----------------------
# Request Models
# -----------------------
class QueryRequest(BaseModel):
    prompt: str
    docs: Optional[List[str]] = []
    mode: Optional[str] = "summarize"
    style: Optional[str] = None


class QuizRequest(BaseModel):
    text: str
    num_questions: Optional[int] = 5


class FeedbackRequest(BaseModel):
    feature: str
    item_name: str
    thumbs_up: int = 0
    rating: int = 0
    note: Optional[str] = ""


# -----------------------
# Routes
# -----------------------
@app.get("/")
def root():
    return {"status": "AI Study Buddy Backend Running ✔️"}


@app.post("/api/upload")
async def upload(file: UploadFile = File(...)):
    """Upload and extract document text"""
    if not file:
        raise HTTPException(status_code=400, detail="No file provided")
    try:
        text = extract_text(file)
        if not text.strip():
            raise HTTPException(status_code=400, detail="Empty or unreadable file.")
        doc_id = f"doc_{len(DOCUMENTS)+1}"
        DOCUMENTS[doc_id] = {"name": file.filename, "text": text}
        return {"id": doc_id, "name": file.filename}
    except Exception as e:
        print("Upload error:", e)
        raise HTTPException(status_code=500, detail=f"Upload failed: {e}")


@app.get("/api/docs")
def docs_list():
    """Return uploaded documents"""
    return [{"id": k, "name": v["name"]} for k, v in DOCUMENTS.items()]


@app.post("/api/query")
def query_handler(q: QueryRequest):
    """Handle both summarization and study planning"""
    try:
        # Check if this is a study planner request (no documents, mode is "chat" or "planner")
        is_planner_request = (
            (q.mode == "chat" or q.mode == "planner") and 
            (not q.docs or len(q.docs) == 0)
        )
        
        if is_planner_request:
            # ===== STUDY PLANNER MODE =====
            print(f"=== Study Planner Request ===")
            print(f"Prompt: {q.prompt}")
            
            # Extract goal and duration from prompt
            duration_match = re.search(r'(\d+)-day', q.prompt)
            duration = int(duration_match.group(1)) if duration_match else 7
            
            goal_match = re.search(r'for learning: (.+?)(?:\n|$)', q.prompt, re.IGNORECASE)
            if not goal_match:
                goal_match = re.search(r'for: (.+?)(?:\n|$)', q.prompt, re.IGNORECASE)
            goal = goal_match.group(1).strip() if goal_match else "Your Topic"
            
            if GEMINI_TEXT_URL:
                try:
                    # Concise prompt for Gemini - forces shorter output
                    concise_prompt = f"""Create a concise {duration}-day study plan for: {goal}

Requirements:
- Be brief and actionable
- Use simple bullet points
- Each day should have:
  * Topic name
  * 3-4 key learning points
  * 1-2 practice tasks
- Keep it under 500 words total
- Use emojis sparingly (only for day markers)
- NO markdown formatting (no **, no ###)

Format:
Day 1: [Topic Name]
• Learn: [Point 1]
• Learn: [Point 2]
• Practice: [Task]

Day 2: [Topic Name]
...

End with 3 quick study tips."""

                    print("Attempting Gemini study plan generation...")
                    answer = call_gemini(concise_prompt)
                    return {"answer": answer.strip(), "sources": []}
                except Exception as e:
                    print(f"Gemini failed for study planner: {e}")
                    # Fall through to local generation
            
            # Local fallback for study planner
            print("Using local study plan generation")
            plan = generate_concise_study_plan(goal, duration)
            return {"answer": plan, "sources": []}
        
        else:
            # ===== REGULAR SUMMARIZATION MODE =====
            selected_docs = [d for d in DOCUMENTS.values() if d["name"] in (q.docs or [])]
            context = "\n\n".join([d["text"][:4000] for d in selected_docs]) or "No context found."

            style_key = (q.style or "simple").lower()
            
            print(f"=== Summarization Request ===")
            print(f"Style: {style_key}")
            print(f"Documents: {[d['name'] for d in selected_docs]}")

            # --- Gemini Preferred ---
            if GEMINI_TEXT_URL:
                try:
                    prompt = get_style_specific_prompt(style_key, context)
                    print(f"Using Gemini with style-specific prompt for: {style_key}")
                    answer = call_gemini(prompt)
                    return {"answer": answer.strip(), "sources": [d["name"] for d in selected_docs]}
                except Exception as e:
                    print(f"Gemini summarization failed, fallback to local: {e}")

            # --- Local Fallback ---
            if selected_docs:
                combined = "\n\n".join([d["text"] for d in selected_docs])
                print(f"Using local fallback for style: {style_key}")
                answer = generate_local_fallback(style_key, combined)
                return {"answer": answer, "sources": [d["name"] for d in selected_docs]}
            else:
                return {"answer": "No documents selected for summarization.", "sources": []}

    except Exception as e:
        print("Query error:", e)
        raise HTTPException(status_code=500, detail=f"Request failed: {e}")


@app.post("/api/quiz")
def quiz(req: QuizRequest):
    """Generate MCQs from text"""
    try:
        if not req.text.strip():
            raise HTTPException(status_code=400, detail="Empty text provided for quiz generation.")

        # Prefer Gemini if available
        if GEMINI_TEXT_URL:
            prompt = f"""
Generate {req.num_questions} multiple choice questions (4 options each) based on the following content.
Return JSON strictly in this format:
[{{"question":"...", "options":["A","B","C","D"], "answer":"A"}}]

Content:
{req.text}
"""
            try:
                result = call_gemini(prompt)
                try:
                    parsed = json.loads(result)
                except Exception:
                    m = re.search(r"(\[.*\])", result, re.S)
                    parsed = json.loads(m.group(1)) if m else []
                return {"quiz": parsed}
            except Exception as e:
                print("Gemini quiz generation failed:", e)

        # Local fallback quiz
        sents = split_sentences(req.text)
        quiz = []
        for i, s in enumerate(sents[:req.num_questions]):
            quiz.append({
                "question": s,
                "options": [s[:60], "Option B", "Option C", "Option D"],
                "answer": "A"
            })
        return {"quiz": quiz}
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Quiz generation failed: {e}")


@app.post("/api/feedback")
def feedback(req: FeedbackRequest):
    """Store user feedback"""
    try:
        FEEDBACKS.append(req.dict())
        return {"message": "Feedback saved successfully", "count": len(FEEDBACKS)}
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Feedback save failed: {e}")


@app.get("/api/feedback")
def feedback_list():
    """Return all feedback"""
    return {"feedbacks": FEEDBACKS}


@app.post("/api/flashcards")
def flashcards(req: QuizRequest):
    """Generate flashcards from text"""
    try:
        if not req.text.strip():
            raise HTTPException(status_code=400, detail="Empty text")

        if GEMINI_TEXT_URL:
            prompt = f"""
Generate {req.num_questions or 10} flashcards from this content.
Return JSON: [{{"front":"Question/Term", "back":"Answer/Definition"}}]

Content:
{req.text}
"""
            try:
                result = call_gemini(prompt)
                parsed = json.loads(result)
                return {"flashcards": parsed}
            except Exception as e:
                print("Gemini flashcard generation failed:", e)

        # Local fallback
        sents = split_sentences(req.text)
        cards = []
        for i, s in enumerate(sents[:10]):
            cards.append({
                "front": f"What about: {s[:50]}...?",
                "back": s
            })
        return {"flashcards": cards}
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Flashcard generation failed: {e}")